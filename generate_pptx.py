#!/usr/bin/env python3
"""Convert markdown outline to PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

# Read markdown file
with open('decks/agent-skills/outline.md', 'r') as f:
    content = f.read()

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Parse slides
lines = content.split('\n')
current_slide = None
current_content = []

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)

# Parse content
slide_title = None
slide_content = []

for line in lines:
    line = line.strip()
    if not line:
        continue
    
    # Presentation title
    if line.startswith('# '):
        continue
    
    # Slide title
    if line.startswith('## Slide'):
        if slide_title and slide_content:
            add_slide(slide_title, slide_content)
        # Extract title after "Slide N: "
        match = re.search(r'Slide \d+: (.+)', line)
        slide_title = match.group(1) if match else line.replace('## ', '')
        slide_content = []
    
    # Bullet points
    elif line.startswith('- '):
        slide_content.append(line[2:].replace('**', '').replace('**', ''))

# Add last slide
if slide_title and slide_content:
    add_slide(slide_title, slide_content)

# Save
output_path = 'decks/agent-skills/agent-skills-deck.pptx'
prs.save(output_path)
print(f"✅ PPTX created: {output_path}")